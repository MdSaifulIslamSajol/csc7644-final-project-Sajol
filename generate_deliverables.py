"""
Script to generate Islam_Saiful_quadchart.pptx and Islam_Saiful_proposal.pdf
for the CSC 7644 Midterm Project.

Run with: conda run -n csc7644 --no-capture-output python generate_deliverables.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib.enums import TA_JUSTIFY, TA_CENTER, TA_LEFT
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, ListFlowable, ListItem
)

# =========================================================================
# PART 1: QUAD CHART (PowerPoint)
# =========================================================================

def create_quad_chart():
    """Create the quad chart by modifying the provided template."""
    template_path = "/home/saiful/Applied_LLM_Development/Midterm project/AINEWSAgentic/Quad Chart Template.pptx"
    output_path = "/home/saiful/Applied_LLM_Development/Midterm project/AINEWSAgentic/Islam_Saiful_quadchart.pptx"

    prs = Presentation(template_path)
    slide = prs.slides[0]

    # Map shape names to content
    content_map = {
        # Title
        "Title 1": "AI News Agentic System: Automated News Curation & Summarization",

        # BLUF
        "TextBox 14": (
            "Bottom Line Up Front: We propose an agentic AI system that autonomously "
            "fetches, summarizes, and delivers curated AI news using a multi-node "
            "LangGraph pipeline with Tavily search integration and Groq-hosted LLMs, "
            "exposed through an interactive Streamlit interface."
        ),

        # Background & Motivation (upper left content)
        "TextBox 15": (
            "• AI news is scattered across hundreds of sources; professionals waste "
            "significant time manually scanning for relevant updates.\n"
            "• Traditional RSS feeds and news aggregators lack intelligent filtering "
            "and summarization capabilities.\n"
            "• LLM-powered agentic systems can automate the full pipeline — search, "
            "filter, summarize, and deliver — reducing information overload.\n"
            "• This project is motivated by the need for a hands-free, reliable "
            "AI news briefing tool for researchers and practitioners."
        ),

        # Data and Resources (upper right content)
        "TextBox 20": (
            "• Tavily Search API: Real-time web search with news-specific filtering "
            "(daily, weekly, monthly time ranges, up to 20 results per query).\n"
            "• Groq Cloud API: Fast LLM inference using Llama 3 (8B, 70B) and "
            "Gemma 2 models for summarization.\n"
            "• All data is fetched live from public web sources via Tavily; no "
            "pre-collected datasets are required.\n"
            "• Output stored as structured Markdown files for easy consumption."
        ),

        # Possible Technical Approach (lower left content)
        "TextBox 18": (
            "• Agentic Architecture: LangGraph-based stateful graph with three "
            "specialized nodes — fetch_news, summarize_news, save_result — "
            "connected via directed edges.\n"
            "• Tool Integration: Tavily search tool bound to chatbot nodes using "
            "LangGraph's ToolNode and tools_condition for conditional routing.\n"
            "• Prompting Strategy: ChatPromptTemplate with structured system prompts "
            "enforcing Markdown output format with dates, summaries, and source URLs.\n"
            "• Three Use Cases: (1) Basic Chatbot, (2) Web-augmented Chatbot with "
            "search tools, (3) AI News pipeline with automated curation.\n"
            "• User Interface: Streamlit web app with sidebar controls for model "
            "selection, API key input, and time-frame filtering."
        ),

        # Impact & References (lower right content)
        "TextBox 19": (
            "• Demonstrates production-ready agentic patterns: tool orchestration, "
            "stateful graph execution, and multi-step reasoning.\n"
            "• Saves 30–60 min/day for AI practitioners needing curated briefings.\n"
            "• Extensible to other domains (finance, healthcare, legal news).\n\n"
            "References:\n"
            "• LangGraph Documentation. (2024). LangChain. https://langchain-ai."
            "github.io/langgraph/\n"
            "• Tavily AI Search API. (2024). https://tavily.com/\n"
            "• Groq Cloud. (2024). https://console.groq.com/"
        ),
    }

    for shape in slide.shapes:
        if shape.name in content_map and shape.has_text_frame:
            tf = shape.text_frame
            new_text = content_map[shape.name]

            if shape.name == "Title 1":
                # Keep title formatting
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.text = ""
                tf.paragraphs[0].runs[0].text = new_text
                tf.paragraphs[0].runs[0].font.size = Pt(28)
                tf.paragraphs[0].runs[0].font.bold = True

            elif shape.name == "TextBox 14":
                # BLUF — bold, slightly smaller
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.text = ""
                tf.paragraphs[0].runs[0].text = new_text
                tf.paragraphs[0].runs[0].font.size = Pt(11)
                tf.paragraphs[0].runs[0].font.bold = True

            else:
                # Content quadrants — clear and rewrite
                tf.clear()
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = new_text
                run.font.size = Pt(9)
                run.font.name = "Calibri"
                p.space_after = Pt(2)
                p.space_before = Pt(2)

    prs.save(output_path)
    print(f"[OK] Quad chart saved: {output_path}")
    return output_path


# =========================================================================
# PART 2: PROPOSAL DOCUMENT (PDF)
# =========================================================================

def create_proposal():
    """Create the 2-3 page proposal document as PDF."""
    output_path = "/home/saiful/Applied_LLM_Development/Midterm project/AINEWSAgentic/Islam_Saiful_proposal.pdf"

    doc = SimpleDocTemplate(
        output_path,
        pagesize=letter,
        topMargin=0.8 * inch,
        bottomMargin=0.8 * inch,
        leftMargin=1 * inch,
        rightMargin=1 * inch,
    )

    styles = getSampleStyleSheet()

    # Custom styles matching the assignment format requirements
    title_style = ParagraphStyle(
        "ProposalTitle",
        parent=styles["Title"],
        fontName="Times-Bold",
        fontSize=14,
        alignment=TA_CENTER,
        spaceAfter=4,
        leading=20,
    )

    subtitle_style = ParagraphStyle(
        "ProposalSubtitle",
        parent=styles["Normal"],
        fontName="Times-Roman",
        fontSize=12,
        alignment=TA_CENTER,
        spaceAfter=10,
        leading=18,
    )

    heading_style = ParagraphStyle(
        "ProposalHeading",
        parent=styles["Heading2"],
        fontName="Times-Bold",
        fontSize=13,
        spaceBefore=10,
        spaceAfter=4,
        leading=18,
    )

    body_style = ParagraphStyle(
        "ProposalBody",
        parent=styles["Normal"],
        fontName="Times-Roman",
        fontSize=12,
        alignment=TA_JUSTIFY,
        leading=18,  # 1.5 line spacing for 12pt
        spaceAfter=6,
    )

    # Build the story
    story = []

    # ---- Title ----
    story.append(Paragraph(
        "AI News Agentic System: Automated News Curation and Summarization "
        "Using LangGraph",
        title_style
    ))
    story.append(Paragraph(
        "Sajol Md Saiful Islam<br/>CSC 7644 — LLM Application Development<br/>Midterm Project Proposal",
        subtitle_style
    ))
    story.append(Spacer(1, 6))

    # ---- 1. Introduction and Motivation ----
    story.append(Paragraph("1. Introduction and Motivation", heading_style))
    story.append(Paragraph(
        "The rapid growth of artificial intelligence research and industry developments has "
        "created an overwhelming volume of daily news. Researchers, engineers, and business "
        "professionals who need to stay current with AI trends face significant information "
        "overload — manually scanning dozens of news sites, blogs, and social media feeds each "
        "day consumes valuable time. Traditional news aggregation tools such as RSS readers "
        "and Google News provide raw collections of articles but lack intelligent filtering, "
        "prioritization, and summarization. Users must still read through numerous headlines "
        "to extract relevant insights.",
        body_style
    ))
    story.append(Paragraph(
        "This project addresses that gap by building an <b>agentic AI system</b> that autonomously "
        "fetches, filters, summarizes, and delivers curated AI news briefings tailored to "
        "user-specified time frames (daily, weekly, or monthly). The primary beneficiaries are "
        "AI practitioners, graduate students, and technology leaders who need concise, reliable "
        "summaries without manual effort. By automating the full news curation pipeline with "
        "LLM-powered agentic workflows, this project demonstrates how prompting, tool "
        "integration, and agentic systems can be combined to solve a practical problem.",
        body_style
    ))

    # ---- 2. Background and Related Work ----
    story.append(Paragraph("2. Background and Related Work", heading_style))
    story.append(Paragraph(
        "Several systems exist for automated news aggregation and summarization. Google News "
        "and Apple News use recommendation algorithms to surface relevant articles, but provide "
        "full articles rather than concise summaries. Tools like Feedly support RSS-based "
        "aggregation with some AI features, but primarily organize content rather than synthesize "
        "it. In the research community, text summarization has been studied from extractive "
        "methods (TextRank, BertSum) to abstractive approaches powered by LLMs.",
        body_style
    ))
    story.append(Paragraph(
        "More recently, <b>agentic AI frameworks</b> such as LangChain and LangGraph have "
        "enabled multi-step workflows where an LLM orchestrates tool calls, maintains state, "
        "and produces structured outputs. Our project combines LangGraph's graph-based agentic "
        "architecture with the Tavily search API for real-time news retrieval and Groq-hosted "
        "LLMs for fast inference, implementing a complete autonomous pipeline orchestrated as "
        "a directed graph with clearly defined node responsibilities.",
        body_style
    ))

    # ---- 3. Proposed Approach ----
    story.append(Paragraph("3. Proposed Approach", heading_style))
    story.append(Paragraph(
        "The project applies two core course concepts: <b>agentic systems with tool use</b> and "
        "<b>advanced prompting strategies</b>. The system is implemented as a LangGraph StateGraph "
        "with three interconnected use cases, each demonstrating different levels of LLM "
        "integration complexity:",
        body_style
    ))
    story.append(Paragraph(
        "<b>Use Case 1 — Basic Chatbot:</b> A single-node graph where user messages flow through "
        "a Groq-hosted LLM (Llama 3 8B/70B or Gemma 2 9B) and return responses. This establishes "
        "the baseline chat interaction pattern.",
        body_style
    ))
    story.append(Paragraph(
        "<b>Use Case 2 — Web-Augmented Chatbot:</b> An enhanced chatbot that binds the Tavily "
        "search tool using LangGraph's ToolNode. The graph uses conditional edges "
        "(tools_condition) to route between the chatbot node and the tool node, enabling the LLM "
        "to autonomously decide when to search the web for additional information before "
        "responding.",
        body_style
    ))
    story.append(Paragraph(
        "<b>Use Case 3 — AI News Pipeline:</b> The primary deliverable — a three-node agentic "
        "pipeline: (1) <i>fetch_news</i> queries Tavily for the top 20 AI news articles within "
        "the selected time range, (2) <i>summarize_news</i> uses a structured ChatPromptTemplate "
        "to instruct the LLM to produce Markdown-formatted summaries sorted by date with source "
        "URLs, and (3) <i>save_result</i> persists the output to a Markdown file for downstream "
        "consumption. The graph executes these nodes sequentially via directed edges.",
        body_style
    ))
    story.append(Paragraph(
        "The user interface is built with Streamlit, providing a sidebar for model selection, "
        "API key management, use case switching, and time-frame controls. This ensures the "
        "system is accessible without command-line expertise.",
        body_style
    ))

    # ---- 4. Data and Resources ----
    story.append(Paragraph("4. Data and Resources", heading_style))
    story.append(Paragraph(
        "The project relies entirely on live data fetched at runtime. <b>Tavily Search API</b> "
        "provides real-time web search with news-specific filtering and configurable time ranges "
        "(daily, weekly, monthly), returning up to 20 results with content snippets, URLs, and "
        "publication dates (free tier: 1,000 searches/month). <b>Groq Cloud API</b> provides "
        "low-latency LLM inference for Llama 3 (8B, 70B) and Gemma 2 (9B) with generous free-tier "
        "rate limits, reducing cost compared to OpenAI while maintaining high summarization quality.",
        body_style
    ))
    story.append(Paragraph(
        "Summarized news is stored as Markdown files in the AINews/ directory "
        "(daily_summary.md, weekly_summary.md, monthly_summary.md), making results easily "
        "viewable and version-controllable. No preprocessing is required since data flows "
        "directly from the Tavily API through the LLM to structured Markdown output.",
        body_style
    ))

    # ---- 5. Evaluation Plan ----
    story.append(Paragraph("5. Evaluation Plan", heading_style))
    story.append(Paragraph(
        "Success is defined along four dimensions: <b>(1) Functional Completeness</b> — all "
        "three use cases must execute end-to-end without errors, with the AI News pipeline "
        "successfully fetching, summarizing, and saving results for all time frames. "
        "<b>(2) Output Quality</b> — summaries will be manually reviewed for factual accuracy, "
        "completeness, and correct Markdown formatting; a sample of 10 summaries will be "
        "spot-checked against source articles.",
        body_style
    ))
    story.append(Paragraph(
        "<b>(3) Latency and Cost</b> — end-to-end execution time will be measured (target: "
        "under 30s for daily, under 60s for monthly), and API costs tracked to confirm free-tier "
        "compliance. <b>(4) Tool Orchestration</b> — the LangGraph conditional routing in the "
        "web-augmented chatbot will be tested with queries that should and should not trigger "
        "web search, verifying correct tools_condition routing.",
        body_style
    ))

    # ---- 6. Feasibility and Scope ----
    story.append(Paragraph("6. Feasibility and Scope", heading_style))
    story.append(Paragraph(
        "This project is designed to be completed within 3–4 weeks. The core architecture "
        "(LangGraph pipeline, Streamlit UI, API integrations) has already been implemented and "
        "tested. The remaining work focuses on evaluation, documentation, and polish.",
        body_style
    ))
    story.append(Paragraph(
        "<b>Risk 1 — API Rate Limits:</b> Both Tavily and Groq impose rate limits on free tiers. "
        "Mitigation: The system fetches only 20 articles per query and uses single-pass "
        "summarization, staying within free-tier quotas. "
        "<b>Risk 2 — LLM Output Inconsistency:</b> Quality may vary across models. Mitigation: "
        "Structured prompts with explicit formatting instructions constrain the output. "
        "<b>Risk 3 — Search Relevance:</b> Results may include irrelevant articles. Mitigation: "
        "Query scoped to \"Top AI technology news\" with topic=\"news\"; domain filtering "
        "can be enabled if needed.",
        body_style
    ))
    story.append(Paragraph(
        "The project scope is intentionally limited to a proof-of-concept demonstrating agentic "
        "workflows, tool orchestration, and LLM-powered summarization — not an enterprise-scale "
        "system, but one that does one thing well with transparent, reproducible architecture.",
        body_style
    ))

    # Build PDF
    doc.build(story)
    print(f"[OK] Proposal saved: {output_path}")
    return output_path


# =========================================================================
# MAIN
# =========================================================================

if __name__ == "__main__":
    print("Generating midterm deliverables...")
    create_quad_chart()
    create_proposal()
    print("Done!")
